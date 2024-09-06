
from PIL import Image
# regular expression
import re 
import pytesseract
# to extract pages from pdf
import fitz 
import io
from pptx import Presentation
from docx import Document



def formatAndCleanText(text):
    formatedText = ' '.join(text.split())
    cleanedText = re.sub(r'[\n\t\r]', '', formatedText)
    return cleanedText



def extractTextFormImage(imageBinary):
    try:
        image = Image.open(io.BytesIO(imageBinary))
        text = pytesseract.image_to_string(image)
        # we arent returning only text because following format is used everywhere
        return ['', formatAndCleanText(text)] 

    except Exception as error:
        print('### error while extracting text from image:\n', error)
        return ['', '']




def extractTextFromPDF(pdfBinary):
    try:
        pdf_document = fitz.open(stream=pdfBinary, filetype="pdf")
        pageNos = len(pdf_document)
        textInPDF = ''
        textInImages = ''

        for page_num in range(pageNos):

            page = pdf_document.load_page(page_num)

            # contains all the images in this page
            image_list = page.get_images(full=True)

            for image in image_list:
                xCoordinate = image[0]
                # contains data of image like width, height, x-co, y-co and other info
                imageData = pdf_document.extract_image(xCoordinate)
                imageBytes = imageData['image']

                #create an image from the binary data
                image = Image.open(io.BytesIO(imageBytes))
                text = pytesseract.image_to_string(image)
                textInImages += text
            
            # for normal text
            textInPDF += page.get_text()

        return [formatAndCleanText(textInPDF), formatAndCleanText(textInImages)]
    
    except Exception as error:
        print('### error while extracting text from pdf:\n', error)
        return ['', '']







def extractTextFromPPTX(pptBinary):
    try:
        presentation = Presentation(io.BytesIO(pptBinary))

        pptText = ''
        textInImages = ''

        for slide in presentation.slides:

            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    pptText += shape.text
                
                elif hasattr(shape, 'image'):
                    imageBlob = io.BytesIO(shape.image.blob)
                    image = Image.open(imageBlob)
                    text = pytesseract.image_to_string(image)
                    textInImages += text
    

        return [formatAndCleanText(pptText), formatAndCleanText(textInImages)]
    
    except Exception as error:
        print('### error while extracting text from ppt:\n', error)
        return ['', '']




def extractTextFromDOCX(docxBinary):
    try:
        doc = Document(io.BytesIO(docxBinary))
        text = ''
        imageText = ''
        
        # extract normal text
        for para in doc.paragraphs:
            text += para.text
        
        # extract text from images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                # Extract the image data as bytes
                image_data = rel.target_part.blob
                imageBlob = io.BytesIO(image_data)
                image = Image.open(imageBlob)
                textInImage = pytesseract.image_to_string(image)
                imageText += textInImage
        text = ""

        return [formatAndCleanText(text), formatAndCleanText(imageText)]
    

    except Exception as error:
        print('### error while extracting text from ppt:\n', error)
        return ['', '']
